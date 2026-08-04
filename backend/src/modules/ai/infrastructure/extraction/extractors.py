"""
Extractores de documentos (PDF, DOCX, PPTX, TXT/MD).

Todos devuelven la misma estructura (`ExtractionResult` con `ContentBlock`),
preservando la ubicación de origen: página o diapositiva. Esa ubicación es lo
que después permite citar "Manual WMS · pág. 23".
"""

from __future__ import annotations

from pathlib import Path

from src.shared.domain.exceptions import DomainError
from src.shared.infrastructure.logging import get_logger

from ...application.ports.extractor import (
    ContentBlock,
    DocumentExtractorPort,
    ExtractionResult,
)

logger = get_logger("ai.extraction")

MIN_CHARS_FOR_TEXT_PAGE = 40  # por debajo de esto se asume página escaneada


class ExtractionFailed(DomainError):
    code = "EXTRACTION_FAILED"
    default_message = "No se pudo extraer el contenido del documento."


class EmptyContent(DomainError):
    code = "EMPTY_CONTENT"
    default_message = "El documento no contiene texto extraíble."


class PdfExtractor(DocumentExtractorPort):
    """PyMuPDF con respaldo OCR para páginas escaneadas."""

    def supports(self, material_type: str) -> bool:
        return material_type == "PDF"

    def extract(self, path: Path) -> ExtractionResult:
        try:
            import fitz  # PyMuPDF
        except ImportError as exc:  # pragma: no cover
            raise ExtractionFailed("PyMuPDF no está instalado.") from exc

        blocks: list[ContentBlock] = []
        used_ocr = False
        order = 0

        try:
            document = fitz.open(str(path))
        except Exception as exc:
            raise ExtractionFailed(f"No se pudo abrir el PDF: {exc}") from exc

        with document:
            for page_number, page in enumerate(document, start=1):
                text = (page.get_text("text") or "").strip()

                if len(text) < MIN_CHARS_FOR_TEXT_PAGE:
                    ocr_text = _ocr_page(page)
                    if ocr_text:
                        text, used_ocr = ocr_text, True

                if not text:
                    continue

                for paragraph in _split_paragraphs(text):
                    blocks.append(ContentBlock(text=paragraph, order=order, page=page_number))
                    order += 1

            page_count = document.page_count

        if not blocks:
            raise EmptyContent(
                "El PDF no tiene texto extraíble. Si es un documento escaneado, "
                "verifique que el OCR esté habilitado."
            )

        return ExtractionResult(blocks=blocks, page_count=page_count, used_ocr=used_ocr)


def _ocr_page(page) -> str:
    """OCR de una página escaneada. Su ausencia no rompe la ingesta."""
    try:
        import pytesseract
        from PIL import Image

        pixmap = page.get_pixmap(dpi=200)
        image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
        return pytesseract.image_to_string(image, lang="spa+eng").strip()
    except Exception:
        logger.debug("OCR no disponible o falló en la página")
        return ""


class DocxExtractor(DocumentExtractorPort):
    """python-docx: conserva la jerarquía de encabezados y las tablas."""

    def supports(self, material_type: str) -> bool:
        return material_type == "DOCX"

    def extract(self, path: Path) -> ExtractionResult:
        try:
            import docx
        except ImportError as exc:  # pragma: no cover
            raise ExtractionFailed("python-docx no está instalado.") from exc

        try:
            document = docx.Document(str(path))
        except Exception as exc:
            raise ExtractionFailed(f"No se pudo abrir el DOCX: {exc}") from exc

        blocks: list[ContentBlock] = []
        current_heading: str | None = None
        order = 0

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            if paragraph.style.name.startswith("Heading"):
                current_heading = text
            blocks.append(ContentBlock(text=text, order=order, heading=current_heading))
            order += 1

        for table in document.tables:
            rows = [
                " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                for row in table.rows
            ]
            content = "\n".join(row for row in rows if row)
            if content:
                blocks.append(
                    ContentBlock(text=f"Tabla:\n{content}", order=order, heading=current_heading)
                )
                order += 1

        if not blocks:
            raise EmptyContent
        return ExtractionResult(blocks=blocks, page_count=0)


class PptxExtractor(DocumentExtractorPort):
    """python-pptx: título, cuerpo y **notas del presentador** por diapositiva."""

    def supports(self, material_type: str) -> bool:
        return material_type == "PPTX"

    def extract(self, path: Path) -> ExtractionResult:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover
            raise ExtractionFailed("python-pptx no está instalado.") from exc

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ExtractionFailed(f"No se pudo abrir el PPTX: {exc}") from exc

        blocks: list[ContentBlock] = []
        order = 0
        slide_count = 0

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_count = slide_number
            title = ""
            body: list[str] = []

            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape == slide.shapes.title:
                    title = text
                else:
                    body.append(text)

            if title:
                blocks.append(
                    ContentBlock(text=title, order=order, page=slide_number, heading=title)
                )
                order += 1

            if body:
                blocks.append(
                    ContentBlock(
                        text="\n".join(body), order=order, page=slide_number, heading=title or None
                    )
                )
                order += 1

            notes = _slide_notes(slide)
            if notes:
                blocks.append(
                    ContentBlock(
                        text=f"Notas del presentador: {notes}",
                        order=order,
                        page=slide_number,
                        heading=title or None,
                    )
                )
                order += 1

        if not blocks:
            raise EmptyContent
        return ExtractionResult(blocks=blocks, page_count=slide_count)


def _slide_notes(slide) -> str:
    try:
        if slide.has_notes_slide:
            return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass
    return ""


class PlainTextExtractor(DocumentExtractorPort):
    """TXT y Markdown, con detección de codificación."""

    def supports(self, material_type: str) -> bool:
        return material_type in {"TXT", "MD"}

    def extract(self, path: Path) -> ExtractionResult:
        raw = path.read_bytes()
        text = _decode(raw)

        blocks = [
            ContentBlock(text=paragraph, order=index)
            for index, paragraph in enumerate(_split_paragraphs(text))
        ]
        if not blocks:
            raise EmptyContent
        return ExtractionResult(blocks=blocks, page_count=0)


def _decode(raw: bytes) -> str:
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(raw).best()
        if best is not None:
            return str(best)
    except Exception:
        pass
    return raw.decode("utf-8", errors="replace")


def _split_paragraphs(text: str) -> list[str]:
    parts = [part.strip() for part in text.split("\n\n")]
    return [part for part in parts if len(part) > 1]
